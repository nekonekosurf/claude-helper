"""
document_reader.py - 各種ファイル形式の読み込みと テキスト抽出

対応形式:
  PDF   - pymupdf (fitz) / pdfplumber / camelot (表抽出)
  PPTX  - python-pptx
  DOCX  - python-docx
  XLSX  - openpyxl
  JPG/PNG - pytesseract (OCR)

依存:
  uv pip install pymupdf pdfplumber python-pptx python-docx openpyxl pytesseract pillow
  apt install tesseract-ocr tesseract-ocr-jpn  # 日本語OCR
"""

from __future__ import annotations

import hashlib
import importlib
import io
import logging
import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


# ──────────────────────────────────────────────────
# データクラス
# ──────────────────────────────────────────────────

@dataclass
class ExtractedDocument:
    """テキスト抽出結果を保持する。"""
    file_path: str
    file_name: str
    file_ext: str
    file_size_bytes: int
    file_hash: str
    text: str                        # 全体テキスト（ページ区切りあり）
    pages: list[str] = field(default_factory=list)   # ページ別テキスト
    tables: list[list[list[str]]] = field(default_factory=list)  # 表データ
    metadata: dict[str, Any] = field(default_factory=dict)       # タイトル・著者等
    extraction_method: str = ""
    error: str | None = None


# ──────────────────────────────────────────────────
# ユーティリティ
# ──────────────────────────────────────────────────

def _sha256(path: str | Path) -> str:
    """ファイルの SHA256 ハッシュを返す。"""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def _try_import(package: str, pip_name: str | None = None):
    """パッケージを try import して、無い場合は None を返す。"""
    try:
        return importlib.import_module(package)
    except ImportError:
        name = pip_name or package
        logger.warning("パッケージが見つかりません: %s  →  uv pip install %s", package, name)
        return None


# ──────────────────────────────────────────────────
# PDF 読み込み
# ──────────────────────────────────────────────────

def read_pdf_pymupdf(path: str | Path) -> ExtractedDocument:
    """
    pymupdf (fitz) で PDF を読み込む。
    高速・日本語対応・デジタル PDF に最適。
    """
    path = Path(path)
    fitz = _try_import("fitz", "pymupdf")
    if fitz is None:
        return ExtractedDocument(
            file_path=str(path), file_name=path.name, file_ext=".pdf",
            file_size_bytes=path.stat().st_size if path.exists() else 0,
            file_hash="", text="", error="pymupdf が未インストールです"
        )

    doc = fitz.open(str(path))
    pages: list[str] = []
    metadata = doc.metadata or {}

    for page in doc:
        pages.append(page.get_text("text"))

    doc.close()
    full_text = "\n\n---ページ区切り---\n\n".join(pages)

    return ExtractedDocument(
        file_path=str(path),
        file_name=path.name,
        file_ext=".pdf",
        file_size_bytes=path.stat().st_size,
        file_hash=_sha256(path),
        text=full_text,
        pages=pages,
        metadata={
            "title":   metadata.get("title", ""),
            "author":  metadata.get("author", ""),
            "subject": metadata.get("subject", ""),
            "creator": metadata.get("creator", ""),
        },
        extraction_method="pymupdf",
    )


def read_pdf_pdfplumber(path: str | Path) -> ExtractedDocument:
    """
    pdfplumber で PDF を読み込む。
    表の抽出精度が高い。複雑なレイアウトに強い。
    """
    path = Path(path)
    pdfplumber = _try_import("pdfplumber")
    if pdfplumber is None:
        return ExtractedDocument(
            file_path=str(path), file_name=path.name, file_ext=".pdf",
            file_size_bytes=0, file_hash="", text="",
            error="pdfplumber が未インストールです"
        )

    pages: list[str] = []
    tables: list[list[list[str]]] = []

    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            # テキスト抽出
            text = page.extract_text() or ""
            pages.append(text)

            # 表抽出
            for table in page.extract_tables():
                if table:
                    # None → 空文字列に変換
                    cleaned = [
                        [cell or "" for cell in row]
                        for row in table
                    ]
                    tables.append(cleaned)

    full_text = "\n\n---ページ区切り---\n\n".join(pages)

    return ExtractedDocument(
        file_path=str(path),
        file_name=path.name,
        file_ext=".pdf",
        file_size_bytes=path.stat().st_size,
        file_hash=_sha256(path),
        text=full_text,
        pages=pages,
        tables=tables,
        extraction_method="pdfplumber",
    )


def read_pdf_with_ocr(path: str | Path, lang: str = "jpn+eng") -> ExtractedDocument:
    """
    スキャン PDF (画像PDF) を OCR で読み込む。
    pymupdf でページを画像に変換 → tesseract でテキスト化。

    前提:
      apt install tesseract-ocr tesseract-ocr-jpn
      uv pip install pymupdf pytesseract pillow
    """
    path = Path(path)
    fitz = _try_import("fitz", "pymupdf")
    pytesseract = _try_import("pytesseract")
    Image_mod = _try_import("PIL.Image", "pillow")

    if not all([fitz, pytesseract, Image_mod]):
        return ExtractedDocument(
            file_path=str(path), file_name=path.name, file_ext=".pdf",
            file_size_bytes=0, file_hash="", text="",
            error="pymupdf / pytesseract / pillow のいずれかが未インストールです"
        )

    Image = Image_mod.Image  # type: ignore[attr-defined]

    doc = fitz.open(str(path))
    pages: list[str] = []

    for page_idx, page in enumerate(doc):
        # 300 dpi 相当で画像に変換
        mat = fitz.Matrix(300 / 72, 300 / 72)
        pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

        try:
            text = pytesseract.image_to_string(img, lang=lang)
        except Exception as e:
            logger.warning("OCR エラー ページ %d: %s", page_idx + 1, e)
            text = ""

        pages.append(text)
        logger.debug("OCR 完了: ページ %d/%d", page_idx + 1, len(doc))

    doc.close()
    full_text = "\n\n---ページ区切り---\n\n".join(pages)

    return ExtractedDocument(
        file_path=str(path),
        file_name=path.name,
        file_ext=".pdf",
        file_size_bytes=path.stat().st_size,
        file_hash=_sha256(path),
        text=full_text,
        pages=pages,
        extraction_method=f"ocr_tesseract_{lang}",
    )


def read_pdf(path: str | Path, force_ocr: bool = False) -> ExtractedDocument:
    """
    PDF 読み込みのメインエントリ。

    1. 通常の digitial PDF → pymupdf で試行
    2. テキストが少ない（スキャンPDF疑い）→ OCR にフォールバック
    3. force_ocr=True → 常に OCR
    """
    path = Path(path)

    if force_ocr:
        logger.info("OCR モードで読み込み: %s", path.name)
        return read_pdf_with_ocr(path)

    result = read_pdf_pymupdf(path)

    # テキストが非常に少ない場合はスキャンPDF と判断して OCR
    text_len = len(result.text.strip())
    if result.error is None and text_len < 100:
        logger.info(
            "テキスト量が少ない (%d文字)。OCR にフォールバック: %s",
            text_len, path.name
        )
        return read_pdf_with_ocr(path)

    return result


# ──────────────────────────────────────────────────
# PowerPoint 読み込み
# ──────────────────────────────────────────────────

def read_pptx(path: str | Path) -> ExtractedDocument:
    """python-pptx で PPTX を読み込む。スライドごとにテキストを抽出。"""
    path = Path(path)
    pptx_mod = _try_import("pptx", "python-pptx")
    if pptx_mod is None:
        return ExtractedDocument(
            file_path=str(path), file_name=path.name, file_ext=".pptx",
            file_size_bytes=0, file_hash="", text="",
            error="python-pptx が未インストールです"
        )

    from pptx import Presentation  # type: ignore[import]
    from pptx.util import Pt       # type: ignore[import]

    prs = Presentation(str(path))
    slides: list[str] = []
    tables: list[list[list[str]]] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = [f"=== スライド {slide_idx} ==="]

        for shape in slide.shapes:
            # テキストフレーム
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_texts.append(line)

            # 表
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
                # テキストとしても追加
                slide_texts.append("[表]")
                for row in table_data:
                    slide_texts.append(" | ".join(row))

        slides.append("\n".join(slide_texts))

    full_text = "\n\n".join(slides)

    return ExtractedDocument(
        file_path=str(path),
        file_name=path.name,
        file_ext=".pptx",
        file_size_bytes=path.stat().st_size,
        file_hash=_sha256(path),
        text=full_text,
        pages=slides,
        tables=tables,
        metadata={"slide_count": len(slides)},
        extraction_method="python-pptx",
    )


# ──────────────────────────────────────────────────
# Word 読み込み
# ──────────────────────────────────────────────────

def read_docx(path: str | Path) -> ExtractedDocument:
    """python-docx で Word ファイルを読み込む。"""
    path = Path(path)
    docx_mod = _try_import("docx", "python-docx")
    if docx_mod is None:
        return ExtractedDocument(
            file_path=str(path), file_name=path.name, file_ext=".docx",
            file_size_bytes=0, file_hash="", text="",
            error="python-docx が未インストールです"
        )

    from docx import Document  # type: ignore[import]
    from docx.oxml.ns import qn  # type: ignore[import]

    doc = Document(str(path))
    paragraphs: list[str] = []
    tables: list[list[list[str]]] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # 見出しスタイルを検出してマーク
            style = para.style.name if para.style else ""
            if "Heading" in style or "見出し" in style:
                text = f"\n## {text}"
            paragraphs.append(text)

    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(row_data)
        tables.append(table_data)

    # コアプロパティ（メタデータ）
    cp = doc.core_properties
    metadata = {
        "title":    cp.title or "",
        "author":   cp.author or "",
        "subject":  cp.subject or "",
        "modified": str(cp.modified) if cp.modified else "",
    }

    full_text = "\n".join(paragraphs)

    return ExtractedDocument(
        file_path=str(path),
        file_name=path.name,
        file_ext=".docx",
        file_size_bytes=path.stat().st_size,
        file_hash=_sha256(path),
        text=full_text,
        tables=tables,
        metadata=metadata,
        extraction_method="python-docx",
    )


# ──────────────────────────────────────────────────
# Excel 読み込み
# ──────────────────────────────────────────────────

def read_xlsx(path: str | Path, max_rows: int = 1000) -> ExtractedDocument:
    """openpyxl で Excel ファイルを読み込む。"""
    path = Path(path)
    openpyxl = _try_import("openpyxl")
    if openpyxl is None:
        return ExtractedDocument(
            file_path=str(path), file_name=path.name, file_ext=".xlsx",
            file_size_bytes=0, file_hash="", text="",
            error="openpyxl が未インストールです"
        )

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sheets: list[str] = []
    tables: list[list[list[str]]] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_texts = [f"=== シート: {sheet_name} ==="]
        sheet_table: list[list[str]] = []
        row_count = 0

        for row in ws.iter_rows(values_only=True):
            if row_count >= max_rows:
                sheet_texts.append(f"... (最大 {max_rows} 行で打ち切り)")
                break
            # None → 空文字に変換
            cells = [str(c) if c is not None else "" for c in row]
            # 完全空行はスキップ
            if any(c for c in cells):
                sheet_table.append(cells)
                sheet_texts.append(" | ".join(cells))
                row_count += 1

        sheets.append("\n".join(sheet_texts))
        if sheet_table:
            tables.append(sheet_table)

    wb.close()
    full_text = "\n\n".join(sheets)

    return ExtractedDocument(
        file_path=str(path),
        file_name=path.name,
        file_ext=".xlsx",
        file_size_bytes=path.stat().st_size,
        file_hash=_sha256(path),
        text=full_text,
        pages=sheets,
        tables=tables,
        metadata={"sheet_count": len(sheets)},
        extraction_method="openpyxl",
    )


# ──────────────────────────────────────────────────
# 画像 OCR
# ──────────────────────────────────────────────────

def read_image_ocr(path: str | Path, lang: str = "jpn+eng") -> ExtractedDocument:
    """画像ファイル (JPG/PNG) を tesseract OCR でテキスト化。"""
    path = Path(path)
    pytesseract = _try_import("pytesseract")
    Image_mod = _try_import("PIL.Image", "pillow")

    if not all([pytesseract, Image_mod]):
        return ExtractedDocument(
            file_path=str(path), file_name=path.name, file_ext=path.suffix,
            file_size_bytes=0, file_hash="", text="",
            error="pytesseract / pillow が未インストールです"
        )

    Image = Image_mod.Image  # type: ignore[attr-defined]
    img = Image.open(str(path))
    text = pytesseract.image_to_string(img, lang=lang)

    return ExtractedDocument(
        file_path=str(path),
        file_name=path.name,
        file_ext=path.suffix,
        file_size_bytes=path.stat().st_size,
        file_hash=_sha256(path),
        text=text,
        pages=[text],
        extraction_method=f"ocr_tesseract_{lang}",
    )


# ──────────────────────────────────────────────────
# 統合エントリポイント
# ──────────────────────────────────────────────────

def read_document(
    path: str | Path,
    force_ocr: bool = False,
) -> ExtractedDocument:
    """
    ファイル拡張子を見て適切なリーダーを自動選択する。

    Args:
        path: ファイルパス
        force_ocr: PDF の場合に強制的に OCR を使う

    Returns:
        ExtractedDocument
    """
    path = Path(path)
    ext = path.suffix.lower()

    dispatch: dict[str, Any] = {
        ".pdf":  lambda p: read_pdf(p, force_ocr=force_ocr),
        ".pptx": read_pptx,
        ".ppt":  read_pptx,
        ".docx": read_docx,
        ".doc":  read_docx,
        ".xlsx": read_xlsx,
        ".xls":  read_xlsx,
        ".jpg":  read_image_ocr,
        ".jpeg": read_image_ocr,
        ".png":  read_image_ocr,
        ".tiff": read_image_ocr,
        ".tif":  read_image_ocr,
        # プレーンテキスト（テスト・メモ用）
        ".txt":  lambda p: ExtractedDocument(
            file_path=str(p), file_name=Path(p).name, file_ext=".txt",
            file_size_bytes=Path(p).stat().st_size,
            file_hash=_sha256(p),
            text=Path(p).read_text(encoding="utf-8", errors="replace"),
            extraction_method="plaintext",
        ),
    }

    reader = dispatch.get(ext)
    if reader is None:
        return ExtractedDocument(
            file_path=str(path),
            file_name=path.name,
            file_ext=ext,
            file_size_bytes=path.stat().st_size if path.exists() else 0,
            file_hash=_sha256(path) if path.exists() else "",
            text="",
            error=f"未対応のファイル形式です: {ext}",
        )

    logger.info("読み込み開始: %s (%s)", path.name, ext)
    return reader(path)


# ──────────────────────────────────────────────────
# 動作確認
# ──────────────────────────────────────────────────
if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("使い方: python document_reader.py <ファイルパス>")
        sys.exit(1)

    result = read_document(sys.argv[1])
    print(f"ファイル名: {result.file_name}")
    print(f"抽出方法: {result.extraction_method}")
    print(f"テキスト長: {len(result.text)} 文字")
    print(f"ページ数: {len(result.pages)}")
    print(f"表の数: {len(result.tables)}")
    if result.error:
        print(f"エラー: {result.error}")
    print("─" * 60)
    print(result.text[:2000])
